from __future__ import annotations

import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
LOCAL_VENDOR = ROOT / ".pptx_build"

if LOCAL_VENDOR.exists():
    sys.path.insert(0, str(LOCAL_VENDOR))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT_DIR = ROOT / "docs"
OUTPUT_PATH = OUTPUT_DIR / "NSS_New_Staff_User_Guide.pptx"
HERO_IMAGE = ROOT / "public" / "assets" / "hero-background.png"
LOGO_IMAGE = ROOT / "public" / "images" / "nss-logo.jpg"

NAVY = RGBColor(15, 23, 42)
SLATE = RGBColor(51, 65, 85)
MUTED = RGBColor(71, 85, 105)
LIGHT_BG = RGBColor(248, 250, 252)
WHITE = RGBColor(255, 255, 255)
CYAN = RGBColor(14, 165, 233)
AMBER = RGBColor(217, 119, 6)
BORDER = RGBColor(203, 213, 225)


SLIDES = [
    {
        "kind": "cover",
        "title": "Nyagatare Secondary School",
        "subtitle": "New Staff User Guide for the NSS Digital System and public website workflows",
        "notes": [
            "Prepared from a full project scan",
            "Updated March 22, 2026",
            "No code included",
        ],
    },
    {
        "title": "System At A Glance",
        "subtitle": "This platform combines public communication with protected school operations.",
        "boxes": [
            {
                "title": "Public Website",
                "body": [
                    "Home page, events, blog/news, board members, enrollment, login, and donation pages.",
                    "Published content and event updates become visible here for visitors, parents, and applicants.",
                ],
            },
            {
                "title": "Secure Staff Workspace",
                "body": [
                    "The /system area is role-aware and only shows modules your account can access.",
                    "Typical staff modules include dashboard, applications, academics, class operations, content, users, finance, invite, and reports.",
                ],
            },
            {
                "title": "Live Records",
                "body": [
                    "Most modules use live Firestore data rather than static demo pages.",
                    "A change in one desk can affect dashboards, student views, applicant tracking, or public pages.",
                ],
            },
        ],
    },
    {
        "title": "Access And First-Time Onboarding",
        "subtitle": "New staff accounts are usually created through the invite workflow before normal sign-in begins.",
        "boxes": [
            {
                "title": "1. Receive An Invite",
                "body": [
                    "A leadership user creates a one-time invite and chooses your role.",
                    "That invite carries the access level you will start with.",
                ],
            },
            {
                "title": "2. Create Your Account",
                "body": [
                    "Open the invite link, confirm your full name, and create your password.",
                    "A valid invite moves you directly into the NSS system after account creation.",
                ],
            },
            {
                "title": "3. Sign In Later",
                "body": [
                    "Go to the Login page, choose Staff, and enter your school email and password.",
                    "Applicants, students, parents, and staff each have separate tabs on the same login screen.",
                ],
            },
            {
                "title": "4. Protect Your Access",
                "body": [
                    "Do not share passwords or leave the system open on shared computers.",
                    "Use My Profile for password changes and sign out from the top bar when finished.",
                ],
            },
        ],
    },
    {
        "title": "Role-Based Access Model",
        "subtitle": "Menus change according to role, so two staff users may not see the same pages.",
        "boxes": [
            {
                "title": "Leadership Roles",
                "body": [
                    "SuperAdmin and Headmaster have the broadest school-wide visibility.",
                    "They can oversee users, roles, reports, admissions, finance, discipline, academics, and publishing.",
                ],
            },
            {
                "title": "Academic Leadership",
                "body": [
                    "DOS, HOD, and DOD focus on academic coordination, class movement, marks visibility, and discipline or content tasks based on role.",
                    "They often bridge leadership reporting with classroom activity.",
                ],
            },
            {
                "title": "Operational Specialists",
                "body": [
                    "Admissions Officer, Bursar, and Content Manager have narrower but important workspaces.",
                    "Their access is centered on admissions, finance, or public communication duties.",
                ],
            },
            {
                "title": "Teaching And Family Accounts",
                "body": [
                    "Teachers mainly work with marks, learning resources, chats, and assigned classes.",
                    "Students and parents mostly view academic information rather than manage school records.",
                ],
            },
        ],
    },
    {
        "title": "Daily Navigation And Account Management",
        "subtitle": "Most staff users begin their day from the dashboard and move into their assigned workspaces.",
        "boxes": [
            {
                "title": "Dashboard Home",
                "body": [
                    "The System Home page is role-aware and shows quick actions, alerts, recent activity, and role-relevant statistics.",
                    "Use it as the safest place to understand what needs attention first.",
                ],
            },
            {
                "title": "My Profile",
                "body": [
                    "Update your own full name, department, and password here.",
                    "Role and sign-in email stay protected from normal self-service editing.",
                ],
            },
            {
                "title": "Users And Roles",
                "body": [
                    "Leadership users can search access profiles, inspect role, department, and status, then update or remove non-protected records.",
                    "The Roles page shows each role and the permissions attached to it.",
                ],
            },
            {
                "title": "Invite Workspace",
                "body": [
                    "Authorized users can create one-time invite links, resend invitation emails, change pending invite roles, or regenerate links.",
                    "Expired or incorrect invites are handled here instead of creating duplicate accounts.",
                ],
            },
        ],
    },
    {
        "title": "Admissions And Applicant Journey",
        "subtitle": "Public applications and internal admissions review are tightly connected.",
        "boxes": [
            {
                "title": "Application Intake",
                "body": [
                    "Staff can open or close public applications from the Admissions Workspace.",
                    "When intake is closed, the enrollment form stops accepting new submissions.",
                ],
            },
            {
                "title": "Review Queue",
                "body": [
                    "Search by applicant name or application ID and filter by status.",
                    "Each record shows applicant details, guardian details, previous school, score, and report links when provided.",
                ],
            },
            {
                "title": "Decision Handling",
                "body": [
                    "Application statuses include Pending, Review, Admitted, Rejected, and Waitlist.",
                    "Staff can save admissions notes, applicant communication notes, and final decision notes.",
                ],
            },
            {
                "title": "Applicant Visibility",
                "body": [
                    "Applicants can track their status in the Applicant Portal using an application ID.",
                    "Decision notes and invite-based applicant account setup can flow back to the applicant side.",
                ],
            },
        ],
    },
    {
        "title": "Class Operations",
        "subtitle": "This desk handles structural class setup before daily teaching begins.",
        "boxes": [
            {
                "title": "Create Classes",
                "body": [
                    "Create a class with name, department, head teacher, subject, academic year, term, and optional student leader.",
                    "This is the foundation for later marks, resources, and class communication.",
                ],
            },
            {
                "title": "Deploy Teachers",
                "body": [
                    "Assign a teacher to a class and subject, then decide whether that assignment allows student invites, parent invites, or class-change action.",
                    "Existing assignments can be edited instead of recreated.",
                ],
            },
            {
                "title": "Place Or Move Students",
                "body": [
                    "Learners can be assigned into a class or moved between classes from one form.",
                    "Movement is logged so class changes remain auditable.",
                ],
            },
            {
                "title": "Leadership Board",
                "body": [
                    "Each class page shows subject coverage, class teacher status, student roster, discipline watch, missing setup items, movement history, and recent actions.",
                    "Use this board before making major class decisions.",
                ],
            },
        ],
    },
    {
        "title": "Academics And Student Support",
        "subtitle": "The Academics workspace connects classroom delivery, marks, messaging, and support visibility.",
        "boxes": [
            {
                "title": "Learning Resources",
                "body": [
                    "Teachers and authorized leaders can publish assignments, exercises, holiday packages, notes, and materials by class and subject.",
                    "Resources can include due dates and attachment links.",
                ],
            },
            {
                "title": "Marks And Comments",
                "body": [
                    "Marks are entered by class, student, subject, term, and academic year.",
                    "Existing mark records update automatically, and teacher comments can be saved with the mark.",
                ],
            },
            {
                "title": "Chats And Follow-Up",
                "body": [
                    "Class Conversations include shared class chats and private subject threads.",
                    "Students can respond to marks and message subject teachers for follow-up.",
                ],
            },
            {
                "title": "What Students See",
                "body": [
                    "The Student Dashboard shows class membership, average mark, learning queue, resources, marks, teacher messaging, and discipline notes.",
                    "Staff should know this view so communication matches what learners can already see.",
                ],
            },
        ],
    },
    {
        "title": "Content And Public Communication",
        "subtitle": "Publishing work in the Content desk directly shapes what the public sees on the website.",
        "boxes": [
            {
                "title": "Stories And Announcements",
                "body": [
                    "Create, search, filter, publish, edit, or delete news, blog posts, and announcements.",
                    "Statuses move through Draft, Review, and Published.",
                ],
            },
            {
                "title": "Event Operations",
                "body": [
                    "Maintain the live school calendar with title, date, category, status, location, and event description.",
                    "Published updates appear on the public Events page.",
                ],
            },
            {
                "title": "Board Members",
                "body": [
                    "Add or update governance profiles for teacher, leader, and parent board categories.",
                    "This information powers the public Board Members page.",
                ],
            },
            {
                "title": "Subscribers",
                "body": [
                    "Newsletter audiences are collected from the homepage footer and Events page.",
                    "Staff can search the list and export subscriber data as CSV when needed.",
                ],
            },
        ],
    },
    {
        "title": "Leadership Operations And Monitoring",
        "subtitle": "Several modules support school-wide oversight beyond day-to-day classroom work.",
        "boxes": [
            {
                "title": "Control Center",
                "body": [
                    "Leadership can review and edit applications, events, donations, and subscriber records from one central page.",
                    "Use it for quick whole-school updates when you do not need to enter each separate desk.",
                ],
            },
            {
                "title": "Finance",
                "body": [
                    "Bursar-focused views track donations, payment status, donor details, references, payment links, and finance health snapshots.",
                    "Donation exports are available for reporting and reconciliation.",
                ],
            },
            {
                "title": "Discipline",
                "body": [
                    "Authorized users monitor discipline cases, severity, ownership, and follow-up status.",
                    "Class operations and student support views also surface discipline context where relevant.",
                ],
            },
            {
                "title": "AI Hub",
                "body": [
                    "Admin users can send operational prompts to the AI Hub for summaries and analytics guidance.",
                    "Non-admin users can still view insight cards even when prompt entry is disabled.",
                ],
            },
        ],
    },
    {
        "title": "Public Pages Every Staff Member Should Know",
        "subtitle": "Even internal users should understand the main public journeys because staff work changes these experiences.",
        "boxes": [
            {
                "title": "Homepage And Login",
                "body": [
                    "The public site introduces the school and links visitors into events, enrollment, donations, and secure login.",
                    "The Login page separates Applicant, Student, Parent, and Staff entry points.",
                ],
            },
            {
                "title": "Enrollment And Applicant Portal",
                "body": [
                    "The Enrollment page collects public applications while applications are open.",
                    "The Applicant Portal lets applicants track status, decision progress, and application details.",
                ],
            },
            {
                "title": "Events And Newsroom",
                "body": [
                    "The Events page displays event details and accepts newsletter subscriptions.",
                    "The Blog page presents published news, blog posts, and announcements in a searchable newsroom format.",
                ],
            },
            {
                "title": "Donations And Governance",
                "body": [
                    "Public donations feed finance records that staff later monitor and verify.",
                    "The Board Members page presents governance profiles that staff maintain in the Content workspace.",
                ],
            },
        ],
    },
    {
        "title": "First-Week Checklist And Common Issues",
        "subtitle": "Use this checklist to help new staff settle in without changing the wrong records.",
        "boxes": [
            {
                "title": "Before Day 1",
                "body": [
                    "Confirm that your invite email, role, and department are correct.",
                    "Know which main workspace you are expected to use: admissions, academics, content, finance, or leadership oversight.",
                ],
            },
            {
                "title": "During First Login",
                "body": [
                    "Open My Profile, set your preferred name and password, and review your visible menu items.",
                    "Visit the dashboard first before editing live records.",
                ],
            },
            {
                "title": "If Data Looks Missing",
                "body": [
                    "Missing menu items usually mean your role does not include that permission.",
                    "Missing student or class data often means the account email, access profile, or class assignment has not been linked yet.",
                ],
            },
            {
                "title": "When To Escalate",
                "body": [
                    "Ask leadership before changing roles, deleting records, publishing public content, or correcting finance data.",
                    "Expired invites, protected accounts, and permission problems should be handled by authorized leadership users.",
                ],
            },
        ],
    },
]


def set_background(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, *, font_size=18, bold=False,
                color=SLATE, font_name="Aptos", align=PP_ALIGN.LEFT) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    run = paragraph.runs[0]
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_body_box(slide, left, top, width, height, title: str, body: list[str]) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(1)

    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, Inches(0.08))
    accent.fill.solid()
    accent.fill.fore_color.rgb = CYAN
    accent.line.fill.background()

    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.18)
    frame.margin_right = Inches(0.18)
    frame.margin_top = Inches(0.14)
    frame.margin_bottom = Inches(0.12)
    frame.vertical_anchor = MSO_ANCHOR.TOP

    title_paragraph = frame.paragraphs[0]
    title_paragraph.text = title
    title_paragraph.alignment = PP_ALIGN.LEFT
    title_run = title_paragraph.runs[0]
    title_run.font.name = "Aptos Display"
    title_run.font.size = Pt(17)
    title_run.font.bold = True
    title_run.font.color.rgb = NAVY
    title_paragraph.space_after = Pt(6)

    for index, line in enumerate(body):
        paragraph = frame.add_paragraph()
        paragraph.text = line
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.level = 0
        paragraph.space_after = Pt(4 if index < len(body) - 1 else 0)
        if paragraph.runs:
            run = paragraph.runs[0]
            run.font.name = "Aptos"
            run.font.size = Pt(11.5)
            run.font.color.rgb = SLATE


def add_slide_chrome(slide, title: str, subtitle: str, slide_number: int) -> None:
    set_background(slide, LIGHT_BG)

    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.88))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()

    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.45), Inches(0.88), Inches(2.0), Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = AMBER
    accent.line.fill.background()

    add_textbox(slide, Inches(0.55), Inches(0.14), Inches(8.8), Inches(0.34), title, font_size=24, bold=True, color=WHITE, font_name="Aptos Display")
    add_textbox(slide, Inches(0.55), Inches(1.02), Inches(11.6), Inches(0.45), subtitle, font_size=11.5, color=MUTED)
    add_textbox(slide, Inches(0.55), Inches(7.0), Inches(4.0), Inches(0.2), "NSS New Staff User Guide", font_size=9, color=MUTED)
    add_textbox(slide, Inches(12.45), Inches(7.0), Inches(0.35), Inches(0.2), str(slide_number), font_size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def add_cover_slide(prs: Presentation, content: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    if HERO_IMAGE.exists():
        slide.shapes.add_picture(str(HERO_IMAGE), 0, 0, width=prs.slide_width, height=prs.slide_height)
    else:
        set_background(slide, NAVY)

    overlay = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = NAVY
    overlay.fill.transparency = 0.18
    overlay.line.fill.background()

    side_panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(0.7), Inches(8.0), Inches(5.6))
    side_panel.fill.solid()
    side_panel.fill.fore_color.rgb = NAVY
    side_panel.fill.transparency = 0.14
    side_panel.line.color.rgb = WHITE
    side_panel.line.transparency = 0.65

    add_textbox(slide, Inches(0.9), Inches(1.12), Inches(6.8), Inches(0.5), content["title"], font_size=28, bold=True, color=WHITE, font_name="Aptos Display")
    add_textbox(slide, Inches(0.9), Inches(1.82), Inches(6.7), Inches(1.2), content["subtitle"], font_size=16, color=WHITE)

    tag = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(0.65), Inches(2.4), Inches(0.38))
    tag.fill.solid()
    tag.fill.fore_color.rgb = CYAN
    tag.line.fill.background()
    add_textbox(slide, Inches(1.05), Inches(0.73), Inches(2.0), Inches(0.16), "Staff onboarding deck", font_size=10, bold=True, color=NAVY)

    note_top = Inches(3.5)
    for note in content["notes"]:
        pill = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.9), note_top, Inches(3.5), Inches(0.52))
        pill.fill.solid()
        pill.fill.fore_color.rgb = WHITE
        pill.fill.transparency = 0.08
        pill.line.fill.background()
        add_textbox(slide, Inches(1.08), note_top + Inches(0.15), Inches(3.1), Inches(0.18), note, font_size=11, bold=True, color=NAVY)
        note_top += Inches(0.68)

    if LOGO_IMAGE.exists():
        slide.shapes.add_picture(str(LOGO_IMAGE), Inches(10.55), Inches(0.9), width=Inches(1.85), height=Inches(1.85))

    add_textbox(slide, Inches(10.08), Inches(3.15), Inches(2.35), Inches(1.0), "NSS", font_size=26, bold=True, color=WHITE, font_name="Aptos Display", align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(9.65), Inches(3.75), Inches(3.2), Inches(0.9), "Public pages, secure staff tools, and everyday role-based workflows", font_size=12, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.9), Inches(6.75), Inches(5.2), Inches(0.2), "Generated from live project structure and user-facing screens", font_size=9.5, color=WHITE)


def box_positions(count: int) -> list[tuple]:
    if count <= 1:
        return [(Inches(0.55), Inches(1.6), Inches(12.2), Inches(4.9))]
    if count == 2:
        return [
            (Inches(0.55), Inches(1.6), Inches(5.98), Inches(4.9)),
            (Inches(6.80), Inches(1.6), Inches(5.98), Inches(4.9)),
        ]
    if count == 3:
        return [
            (Inches(0.55), Inches(1.75), Inches(3.92), Inches(4.65)),
            (Inches(4.70), Inches(1.75), Inches(3.92), Inches(4.65)),
            (Inches(8.85), Inches(1.75), Inches(3.92), Inches(4.65)),
        ]
    return [
        (Inches(0.55), Inches(1.78), Inches(5.98), Inches(2.2)),
        (Inches(6.80), Inches(1.78), Inches(5.98), Inches(2.2)),
        (Inches(0.55), Inches(4.15), Inches(5.98), Inches(2.2)),
        (Inches(6.80), Inches(4.15), Inches(5.98), Inches(2.2)),
    ]


def add_content_slide(prs: Presentation, content: dict, slide_number: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide, content["title"], content["subtitle"], slide_number)

    for position, box in zip(box_positions(len(content["boxes"])), content["boxes"]):
        add_body_box(slide, *position, box["title"], box["body"])


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_cover_slide(prs, SLIDES[0])
    for slide_number, content in enumerate(SLIDES[1:], start=2):
        add_content_slide(prs, content, slide_number)
    return prs


def main() -> None:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    presentation = build_presentation()
    presentation.save(str(OUTPUT_PATH))
    print(OUTPUT_PATH)


if __name__ == "__main__":
    main()
